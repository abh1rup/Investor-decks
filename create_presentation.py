#!/usr/bin/env python3
"""
ChatterBot Investor Presentation Generator
Generates a professional PowerPoint deck for investor pitches
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Color scheme
PRIMARY_COLOR = RGBColor(52, 73, 94)      # Dark slate blue
ACCENT_COLOR = RGBColor(41, 182, 246)    # Bright blue
SUCCESS_COLOR = RGBColor(76, 175, 80)    # Green
WARNING_COLOR = RGBColor(255, 152, 0)    # Orange
TEXT_COLOR = RGBColor(33, 33, 33)        # Dark gray
LIGHT_BG = RGBColor(240, 248, 255)       # Alice blue

def create_title_slide(prs):
    """Create the opening title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "ChatterBot"
    title_p.font.size = Pt(72)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "AI-Powered Conversational Intelligence Platform"
    subtitle_p.font.size = Pt(28)
    subtitle_p.font.color.rgb = ACCENT_COLOR
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    tagline_frame = tagline_box.text_frame
    tagline_p = tagline_frame.paragraphs[0]
    tagline_p.text = "Transform Customer Engagement with Enterprise-Grade AI"
    tagline_p.font.size = Pt(18)
    tagline_p.font.italic = True
    tagline_p.font.color.rgb = RGBColor(200, 200, 200)
    tagline_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_points, slide_number=None):
    """Add a standard content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = PRIMARY_COLOR
    header_shape.line.color.rgb = PRIMARY_COLOR
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Slide number
    if slide_number:
        num_box = slide.shapes.add_textbox(Inches(9), Inches(0.15), Inches(0.8), Inches(0.5))
        num_frame = num_box.text_frame
        num_p = num_frame.paragraphs[0]
        num_p.text = str(slide_number)
        num_p.font.size = Pt(16)
        num_p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.2))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(content_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.line_spacing = 1.2

def create_presentation():
    """Create the complete investor presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    create_title_slide(prs)
    
    # Slide 2: Problem Statement
    add_content_slide(prs, "The Problem", [
        "• Enterprises struggle with customer engagement at scale",
        "• Traditional chatbots lack intelligence and personalization",
        "• Integration with existing systems is complex and costly",
        "• High maintenance overhead and poor user experiences",
        "• Market demand for AI-powered conversational solutions growing 45% YoY"
    ], 2)
    
    # Slide 3: Our Solution
    add_content_slide(prs, "ChatterBot Solution", [
        "• Enterprise-grade AI conversational platform",
        "• Easy integration with existing systems via REST API",
        "• Real-time chat with context-aware responses",
        "• Production-ready with security & compliance features",
        "• Scalable architecture supporting millions of conversations"
    ], 3)
    
    # Slide 4: Key Features
    add_content_slide(prs, "Key Features & Capabilities", [
        "✓ AI-Powered Intelligence: Advanced NLP and ML models",
        "✓ Real-Time Messaging: Instant, reliable message delivery",
        "✓ Security-First: Enterprise-grade encryption & authentication",
        "✓ Scalable Architecture: Handle millions of concurrent users",
        "✓ Mobile-Optimized: Seamless experience across all devices",
        "✓ Analytics Dashboard: Real-time insights and metrics"
    ], 4)
    
    # Slide 5: Technology Stack
    add_content_slide(prs, "Modern Technology Stack", [
        "Frontend: React 18 with Vite (40% faster builds)",
        "Backend: Python with production-grade security",
        "Database: PostgreSQL for reliability & consistency",
        "Caching: Redis for sub-second response times",
        "Infrastructure: Docker containerization with K8s ready",
        "CI/CD: GitHub Actions for automated testing & deployment"
    ], 5)
    
    # Slide 6: Competitive Advantages
    add_content_slide(prs, "Why ChatterBot Wins", [
        "🔒 Security: Rate limiting, production-grade validation",
        "⚡ Performance: 40% smaller bundle, code splitting optimization",
        "🛠️ Developer Friendly: Well-documented APIs, easy integration",
        "📊 Enterprise Ready: 70% code coverage, automated testing",
        "🚀 Scalable: Cloud-native architecture, auto-scaling support",
        "💰 Cost Effective: Open-source with commercial support"
    ], 6)
    
    # Slide 7: Enterprise Features
    add_content_slide(prs, "Enterprise-Grade Quality", [
        "Security: Rate limiting (5/min auth, 100/min API)",
        "Quality Assurance: 70% minimum code coverage required",
        "Automated Testing: Unit, integration, and E2E tests",
        "Security Scanning: Weekly dependency & vulnerability audits",
        "Monitoring: Health checks & real-time alerting",
        "Compliance: GDPR-ready with data retention policies"
    ], 7)
    
    # Slide 8: Market Opportunity
    add_content_slide(prs, "Market Opportunity", [
        "Global Conversational AI Market: $15.8B (2023)",
        "Projected CAGR: 23.5% through 2030",
        "Target TAM: Enterprise customer service ($8.2B segment)",
        "Early adopter advantage in AI-powered platforms",
        "Expansion opportunities: E-commerce, Healthcare, Finance",
        "Average contract value: $50K-500K annually"
    ], 8)
    
    # Slide 9: Use Cases
    add_content_slide(prs, "Real-World Applications", [
        "Customer Support: 24/7 intelligent support automation",
        "Sales Enablement: Lead qualification & engagement",
        "HR Automation: Employee onboarding & FAQs",
        "Internal Tools: IT helpdesk automation",
        "E-commerce: Product recommendations & customer service",
        "Healthcare: Patient engagement & appointment scheduling"
    ], 9)
    
    # Slide 10: Product Roadmap
    add_content_slide(prs, "Product Roadmap", [
        "Q3 2026: Multi-language support & advanced NLP models",
        "Q4 2026: Mobile app & offline capabilities",
        "Q1 2027: Industry-specific templates (Finance, Healthcare, Retail)",
        "Q2 2027: Advanced analytics & predictive insights",
        "Q3 2027: WhatsApp & SMS channel integrations",
        "Q4 2027: Enterprise marketplace for plugins & extensions"
    ], 10)
    
    # Slide 11: Financial Projections
    add_content_slide(prs, "Financial Projections", [
        "Year 1: 100 enterprise customers, $5M ARR",
        "Year 2: 350 customers, $18M ARR (260% growth)",
        "Year 3: 800 customers, $45M ARR (150% growth)",
        "Gross Margin: 75-80% (software business model)",
        "CAC Payback: 8-10 months",
        "Expansion revenue from upsells & add-ons: 30% of MRR"
    ], 11)
    
    # Slide 12: Team & Expertise
    add_content_slide(prs, "Team & Leadership", [
        "Founder & CEO: Full-stack expertise in AI & SaaS",
        "CTO: 10+ years in distributed systems & cloud architecture",
        "VP Sales: Track record scaling B2B SaaS companies",
        "Head of Product: Former PM at leading AI platform",
        "Backing: [Investors/Advisors - customize as needed]",
        "Advisory Board: Industry experts from Fortune 500 companies"
    ], 12)
    
    # Slide 13: Funding Requirements
    add_content_slide(prs, "Funding & Use of Proceeds", [
        "Seeking: $2M Series A funding",
        "• 40% - Product Development & AI/ML talent",
        "• 35% - Sales & Marketing team expansion",
        "• 15% - Infrastructure & DevOps",
        "• 10% - Operations & Finance",
        "Expected Timeline: 18-24 months to Series B"
    ], 13)
    
    # Slide 14: Call to Action
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Main message
    cta_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    cta_frame = cta_box.text_frame
    cta_frame.word_wrap = True
    cta_p = cta_frame.paragraphs[0]
    cta_p.text = "Join Us in Transforming Customer Engagement"
    cta_p.font.size = Pt(44)
    cta_p.font.bold = True
    cta_p.font.color.rgb = RGBColor(255, 255, 255)
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_p = contact_frame.paragraphs[0]
    contact_p.text = "Let's Talk\nabh1rup@example.com | (555) 123-4567"
    contact_p.font.size = Pt(24)
    contact_p.font.color.rgb = ACCENT_COLOR
    contact_p.alignment = PP_ALIGN.CENTER
    contact_p.line_spacing = 1.5
    
    # Save presentation
    filename = "ChatterBot-Investor-Deck.pptx"
    prs.save(filename)
    print(f"✅ Presentation created: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print(f"📅 Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

if __name__ == "__main__":
    create_presentation()
